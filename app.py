"""
app.py — Streamlit Web App
AI-Powered Document Summarization System
TEYZIX CORE Internship | Task AI-INT-1 | June Batch 2026
"""

import streamlit as st
import string
from collections import Counter
import math
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
import io

# ── NLTK Data ─────────────────────────────────────────────
for pkg in ["punkt", "punkt_tab", "stopwords"]:
    try:
        nltk.data.find(f"tokenizers/{pkg}" if "punkt" in pkg else f"corpora/{pkg}")
    except LookupError:
        nltk.download(pkg, quiet=True)

# ── Page Config ───────────────────────────────────────────
st.set_page_config(
    page_title="AI Document Summarizer — TEYZIX CORE",
    page_icon="🤖",
    layout="wide"
)

# ── Custom CSS ────────────────────────────────────────────
st.markdown("""
<style>
    .main { background-color: #0E1117; }
    .stTextArea textarea { font-size: 14px; }
    .keyword-tag {
        display: inline-block;
        background: #1D4ED8;
        color: white;
        padding: 3px 10px;
        border-radius: 20px;
        margin: 3px;
        font-size: 13px;
    }
    .header-banner {
        background: linear-gradient(135deg, #0F6E56, #1D4ED8);
        padding: 20px;
        border-radius: 12px;
        text-align: center;
        margin-bottom: 25px;
    }
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════
# FILE READERS
# ══════════════════════════════════════════════════════════

def read_txt(file):
    return file.read().decode("utf-8", errors="ignore")

def read_pdf(file):
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
        return text.strip()
    except Exception as e:
        return None

def read_docx(file):
    try:
        import docx
        doc = docx.Document(file)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        return None

def read_pptx(file):
    try:
        from pptx import Presentation
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return None


# ══════════════════════════════════════════════════════════
# NLP FUNCTIONS
# ══════════════════════════════════════════════════════════

def preprocess(text):
    sentences = sent_tokenize(text)
    tokens = word_tokenize(text.lower())
    stop_words = set(stopwords.words("english"))
    filtered = [t for t in tokens if t not in string.punctuation
                and t not in stop_words and not t.isdigit()]
    return sentences, filtered

def word_freq(filtered_tokens):
    return Counter(filtered_tokens)

def tfidf_scores(sentences):
    stop_words = set(stopwords.words("english"))
    N = len(sentences)
    bags = []
    for s in sentences:
        tokens = word_tokenize(s.lower())
        bag = Counter([t for t in tokens if t not in string.punctuation
                       and t not in stop_words and not t.isdigit()])
        bags.append(bag)
    df = Counter()
    for bag in bags:
        for w in bag:
            df[w] += 1
    idf = {w: math.log((N + 1) / (c + 1)) + 1 for w, c in df.items()}
    scores = {}
    for i, (s, bag) in enumerate(zip(sentences, bags)):
        total = sum(bag.values()) or 1
        scores[i] = sum((c / total) * idf.get(w, 0) for w, c in bag.items())
    return scores

def freq_scores(sentences, freq_map):
    max_f = max(freq_map.values(), default=1)
    norm = {w: f / max_f for w, f in freq_map.items()}
    stop_words = set(stopwords.words("english"))
    scores = {}
    for i, s in enumerate(sentences):
        tokens = word_tokenize(s.lower())
        scores[i] = sum(norm.get(t, 0) for t in tokens
                        if t not in string.punctuation and t not in stop_words)
    return scores

def summarize(text, n, method):
    sentences, filtered = preprocess(text)
    freq = word_freq(filtered)
    n = max(1, min(n, len(sentences)))
    if method == "TF-IDF":
        scores = tfidf_scores(sentences)
    elif method == "Frequency":
        scores = freq_scores(sentences, freq)
    else:
        def norm(d):
            mx = max(d.values(), default=1) or 1
            return {k: v / mx for k, v in d.items()}
        f = norm(freq_scores(sentences, freq))
        t = norm(tfidf_scores(sentences))
        scores = {i: (f.get(i, 0) + t.get(i, 0)) / 2 for i in range(len(sentences))}
    top = sorted(sorted(scores, key=scores.get, reverse=True)[:n])
    summary = " ".join(sentences[i] for i in top)
    return summary, freq, scores, sentences


# ══════════════════════════════════════════════════════════
# UI
# ══════════════════════════════════════════════════════════

st.markdown("""
<div class="header-banner">
    <h1 style="color:white; margin:0;">🤖 AI Document Summarization System</h1>
    <p style="color:#A0AEC0; margin:5px 0 0 0;">TEYZIX CORE Internship | Task AI-INT-1 | June Batch 2026</p>
</div>
""", unsafe_allow_html=True)

with st.sidebar:
    st.image("https://img.shields.io/badge/TEYZIX-CORE-green?style=for-the-badge")
    st.markdown("### ⚙️ Settings")
    method = st.selectbox("Summarization Method", ["TF-IDF", "Frequency", "Combined"])
    num_sentences = st.slider("Number of sentences in summary", 1, 10, 3)
    st.markdown("---")
    st.markdown("### 📖 About")
    st.markdown("""
    This AI system uses **NLP** to extract the most important sentences from any document.
    
    **Supports:**
    - 📝 Plain Text
    - 📄 PDF Files
    - 📘 Word Documents (.docx)
    - 📊 PowerPoint (.pptx)
    
    **Built with:** Python + NLTK
    """)
    st.markdown("---")
    st.markdown("**👨‍💻 Saqlain**  \nAI Intern — TEYZIX CORE")

tab1, tab2 = st.tabs(["📝 Text Input", "📁 File Upload (PDF/DOCX/PPTX/TXT)"])

def show_results(text, method, num_sentences):
    if len(text.split()) < 30:
        st.warning("⚠️ Text is too short. Please provide a longer document.")
        return
    with st.spinner("🧠 Analyzing document..."):
        summary, freq, scores, sentences = summarize(text, num_sentences, method)
    orig_words = len(text.split())
    summ_words = len(summary.split())
    compression = round(summ_words / orig_words * 100, 1)
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Original Words", orig_words)
    c2.metric("Summary Words", summ_words)
    c3.metric("Compression", f"{compression}%")
    c4.metric("Method", method)
    st.markdown("### ✅ Summary")
    st.success(summary)
    st.markdown("### 🔑 Top Keywords")
    top_kw = freq.most_common(10)
    kw_html = " ".join(f'<span class="keyword-tag">{w} ({c})</span>' for w, c in top_kw)
    st.markdown(kw_html, unsafe_allow_html=True)
    col_a, col_b = st.columns(2)
    with col_a:
        st.markdown("**Original Text**")
        st.text_area("", text[:2000] + ("..." if len(text) > 2000 else ""), height=200, disabled=True, label_visibility="collapsed")
    with col_b:
        st.markdown("**Summary**")
        st.text_area("", summary, height=200, disabled=True, label_visibility="collapsed")
    export_text = f"SUMMARY\n{'='*50}\n{summary}\n\nTOP KEYWORDS\n{'='*50}\n" + "\n".join(f"{w}: {c}" for w, c in top_kw)
    st.download_button("⬇️ Download Summary (.txt)", export_text, file_name="summary_output.txt", mime="text/plain")

with tab1:
    text_input = st.text_area("Paste your document here:", height=200,
        placeholder="Paste any long article, report, or document here...")
    if st.button("🚀 Summarize", type="primary"):
        if not text_input.strip():
            st.error("⚠️ Please paste some text first!")
        else:
            show_results(text_input, method, num_sentences)

with tab2:
    st.markdown("**Supported formats: TXT, PDF, DOCX, PPTX**")
    uploaded = st.file_uploader("Upload your document", type=["txt", "pdf", "docx", "pptx"])
    if uploaded:
        ext = uploaded.name.split(".")[-1].lower()
        text_file = None
        with st.spinner(f"📂 Reading {ext.upper()} file..."):
            if ext == "txt":
                text_file = read_txt(uploaded)
            elif ext == "pdf":
                text_file = read_pdf(uploaded)
            elif ext == "docx":
                text_file = read_docx(uploaded)
            elif ext == "pptx":
                text_file = read_pptx(uploaded)

        if not text_file:
            st.error("❌ Could not read file. Make sure it has readable text (not scanned image).")
        else:
            st.success(f"✅ File loaded: {uploaded.name} | {len(text_file.split())} words found")
            if st.button("🚀 Summarize File", type="primary"):
                show_results(text_file, method, num_sentences)

st.markdown("---")
st.markdown("<p style='text-align:center; color:#4A5568;'>Built by Saqlain | TEYZIX CORE AI Internship | June 2026 | Task AI-INT-1</p>", unsafe_allow_html=True)
